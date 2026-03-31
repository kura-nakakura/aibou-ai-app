import urllib.parse
import json
import re
import io
import base64
from gtts import gTTS
from streamlit_mic_recorder import speech_to_text
import google.generativeai as genai
import os
import streamlit as st

# 🚨 パワポ生成用のライブラリ
try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    st.error("⚠️ `python-pptx` ライブラリがインストールされていません。requirements.txt を確認してください。")

/* 🚀 新機能：ホバー時に下部に表示される説明エリア */
    .desc-display-area {
        height: 80px;
        border: 1px solid rgba(255, 255, 255, 0.1);
        border-radius: 12px;
        margin-top: 40px;
        display: flex;
        justify-content: center;
        align-items: center;
        background: rgba(10, 20, 40, 0.4);
        transition: all 0.3s ease;
        text-align: center;
        padding: 0 20px;
    }
    .default-desc { color: #718096; font-weight: bold; letter-spacing: 2px; }
    .app-desc, .img-desc, .vid-desc, .slide-desc {
        display: none;
        color: #ffffff;
        font-size: 14px;
        font-weight: bold;
        letter-spacing: 1px;
        line-height: 1.5;
        text-shadow: 0 0 10px rgba(255, 255, 255, 0.5), 0 0 20px rgba(0, 243, 255, 0.3);
    }

    /* いずれかのボタンがホバーされたらデフォルトテキストを消し、枠を光らせる */
    [data-testid="stVerticalBlock"]:has(button:hover) .default-desc { display: none; }
    [data-testid="stVerticalBlock"]:has(button:hover) .desc-display-area {
        border-color: #ffffff;
        box-shadow: 0 0 20px rgba(255, 255, 255, 0.4), inset 0 0 10px rgba(0, 243, 255, 0.2);
        background: rgba(20, 35, 60, 0.8);
    }

    /* 各ボタンのホバーに連動して対応する説明を表示する魔法のCSS */
    [data-testid="stVerticalBlock"]:has([data-testid="column"]:has(.hover-target-app) button:hover) .app-desc { display: block; }
    [data-testid="stVerticalBlock"]:has([data-testid="column"]:has(.hover-target-img) button:hover) .img-desc { display: block; }
    [data-testid="stVerticalBlock"]:has([data-testid="column"]:has(.hover-target-vid) button:hover) .vid-desc { display: block; }
    [data-testid="stVerticalBlock"]:has([data-testid="column"]:has(.hover-target-slide) button:hover) .slide-desc { display: block; }

# 💎 UIデザイン用CSS (ボタンの白飛び解消 ＆ 白ネオン反応)
st.markdown("""
    <style>
    /* 1. 全体をダーク＆サイバーパンクな雰囲気に */
    [data-testid="stAppViewContainer"] {
        background-color: #030b14 !important;
        background-image: radial-gradient(circle at 50% 120%, rgba(0, 150, 255, 0.15), transparent) !important;
        color: #e2e8f0 !important;
    }

    /* 🌟 2. ボス指定の完全再現: シアンのロゴ */
    .central-logo {
        text-align: center;
        font-size: 70px;
        color: #00f3ff !important; /* 絶対に白にしない。完全なシアン */
        text-shadow: 0 0 15px rgba(0, 243, 255, 0.8), 0 0 30px rgba(0, 150, 255, 0.5) !important;
        margin: 30px 0 10px 0;
    }

    /* 🌟 3. ボス指定の完全再現: タイトル(白文字 + シアン光彩) */
    .saas-title {
        color: #ffffff !important;
        font-weight: 900;
        letter-spacing: 6px;
        margin-bottom: 5px;
        text-shadow: 0 0 10px rgba(0, 243, 255, 0.6), 0 0 20px rgba(0, 150, 255, 0.4) !important;
        text-align: center;
    }
    .central-logo-sub {
        text-align: center;
        color: #a0aec0;
        font-size: 12px;
        letter-spacing: 4px;
        margin-top: -10px;
        margin-bottom: 40px;
    }

    /* 🚀 3. 【修正】ボタンの白飛びを解消し、ダークサイバー仕様に変更 */
    div.stButton > button {
        background-color: rgba(15, 23, 42, 0.8) !important; /* ボタンの背景をダークにして文字を見せる */
        border: 1px solid rgba(255, 255, 255, 0.3) !important; /* 普段は静かな白い枠線 */
        border-radius: 12px !important;
        color: #ffffff !important; /* 文字は真っ白 */
        transition: all 0.3s ease !important;
        padding: 10px !important;
    }
    div.stButton > button p {
        color: #ffffff !important;
        font-weight: 800 !important;
        letter-spacing: 2px !important;
    }

    /* 🔥 4. ホバー時のエフェクト (ボスの要望：白いネオンの反応！) */
    div.stButton > button:hover {
        background-color: rgba(30, 41, 59, 0.9) !important;
        border-color: #ffffff !important; /* 枠線が白く発光 */
        box-shadow: 0 0 20px rgba(255, 255, 255, 0.6), inset 0 0 10px rgba(255, 255, 255, 0.2) !important;
        transform: translateY(-3px) !important; /* 少しだけ浮く */
    }
    div.stButton > button:hover p {
        text-shadow: 0 0 10px #ffffff, 0 0 20px #ffffff !important; /* 文字が真っ白なネオンに */
    }
    
    /* 下部の光るプラットフォーム */
    .hologram-platform {
        position: fixed;
        bottom: -80px;
        left: 50%;
        transform: translateX(-50%);
        width: 600px;
        height: 150px;
        background: radial-gradient(ellipse at 50% 50%, rgba(0, 243, 255, 0.15), transparent 70%);
        border-radius: 50%;
        box-shadow: 0 0 50px rgba(0, 243, 255, 0.2);
        pointer-events: none;
        z-index: -1;
    }
    
    /* サブ画面のコンテナ */
    [data-testid="stVerticalBlockBorderWrapper"] {
        background: rgba(15, 23, 42, 0.6) !important;
        backdrop-filter: blur(12px) !important;
        border: 1px solid #2d3748 !important;
        border-radius: 15px !important;
    }
    </style>
""", unsafe_allow_html=True)

# 🚨 初期化
if "auto_fix_prompt" not in st.session_state: st.session_state.auto_fix_prompt = ""
if "forge_workspaces" not in st.session_state: st.session_state.forge_workspaces = {}
if "current_forge_ws" not in st.session_state: st.session_state.current_forge_ws = None 
if "ai_voice_base64" not in st.session_state: st.session_state.ai_voice_base64 = None
if "just_generated_audio" not in st.session_state: st.session_state.just_generated_audio = False
if "selected_forge_mode" not in st.session_state: st.session_state.selected_forge_mode = None

# ==========================================
# 🚪 ステージ1：ホログラムカード選択画面 (横4列 ＋ ホバー説明機能)
# ==========================================
if st.session_state.current_forge_ws is None and st.session_state.selected_forge_mode is None:
    st.markdown('<div class="central-logo">⬡</div>', unsafe_allow_html=True)
    st.markdown("<h2 class='saas-title'>❖ FORGE STUDIO ❖</h2>", unsafe_allow_html=True)
    st.markdown("<p class='central-logo-sub'>SELECT SYSTEM ENGINE</p>", unsafe_allow_html=True)
    
    st.markdown('<div class="hologram-platform"></div>', unsafe_allow_html=True)

    # 4つのボタンと、それをCSSで検知するための「見えないマーカー」を設置
    c1, c2, c3, c4 = st.columns(4, gap="large")
    with c1:
        st.markdown('<div class="hover-target-app" style="display:none;"></div>', unsafe_allow_html=True)
        if st.button("APP STUDIO", use_container_width=True): 
            st.session_state.selected_forge_mode = "APP"; st.rerun()
    with c2:
        st.markdown('<div class="hover-target-img" style="display:none;"></div>', unsafe_allow_html=True)
        if st.button("IMAGE GENERATOR", use_container_width=True): 
            st.session_state.selected_forge_mode = "IMAGE"; st.rerun()
    with c3:
        st.markdown('<div class="hover-target-vid" style="display:none;"></div>', unsafe_allow_html=True)
        if st.button("VIDEO PRODUCTION", use_container_width=True): 
            st.session_state.selected_forge_mode = "VIDEO"; st.rerun()
    with c4:
        st.markdown('<div class="hover-target-slide" style="display:none;"></div>', unsafe_allow_html=True)
        if st.button("SLIDE DECK", use_container_width=True): 
            st.session_state.selected_forge_mode = "SLIDE"; st.rerun()

    # 🚀 4つのボタンの下に表示される動的説明エリア
    st.markdown("""
        <div class="desc-display-area">
            <span class="default-desc">HOVER OVER AN ENGINE TO VIEW SPECIFICATIONS</span>
            <span class="app-desc">🤖 <b>[ APP STUDIO ]</b><br>ボスの指示から、美しくバグのないアプリケーションのUIとロジックを自律的に構築します。</span>
            <span class="img-desc">🎨 <b>[ IMAGE GENERATOR ]</b><br>画像生成AIのための完璧な英語プロンプトを構築し、最高の1枚を引き出します。</span>
            <span class="vid-desc">🎬 <b>[ VIDEO PRODUCTION ]</b><br>最先端動画生成AIに向けた、プロ品質の絵コンテとカメラワーク指定を作成します。</span>
            <span class="slide-desc">📊 <b>[ SLIDE DECK ]</b><br>論理的なプレゼン構成を考案し、説得力のあるスライド資料(.pptx)を即座に出力します。</span>
        </div>
    """, unsafe_allow_html=True)
# ==========================================
# 🚪 ステージ2：モード別プロジェクト管理画面
# ==========================================
elif st.session_state.current_forge_ws is None and st.session_state.selected_forge_mode is not None:
    mode = st.session_state.selected_forge_mode
    
    col_back, col_title = st.columns([2, 8])
    with col_back:
        if st.button("← BACK TO ENGINES", use_container_width=True):
            st.session_state.selected_forge_mode = None
            st.rerun()
    with col_title:
        st.markdown(f"<h3 style='color:#ffffff; font-weight:800; margin-top:-5px;'>[ {mode} ENGINE ]</h3>", unsafe_allow_html=True)
    
    st.markdown("---")
    
    col_create, col_list = st.columns([4, 6], gap="large")
    
    with col_create:
        st.markdown("<p style='font-weight:bold; color:#a0aec0; font-size:12px;'>[ INITIALIZE NEW PROJECT ]</p>", unsafe_allow_html=True)
        with st.container(border=True):
            new_ws_name = st.text_input("PROJECT NAME", label_visibility="collapsed", placeholder="Enter project name...")
            st.markdown("<br>", unsafe_allow_html=True)
            if st.button("CREATE WORKSPACE", type="primary", use_container_width=True):
                if new_ws_name and new_ws_name not in st.session_state.forge_workspaces:
                    st.session_state.forge_workspaces[new_ws_name] = {
                        "type": mode, "chat": [], "code": "", "media": "", "retries": 0
                    }
                    st.session_state.current_forge_ws = new_ws_name
                    st.rerun()
                elif not new_ws_name:
                    st.error("Please enter a project name.")

    with col_list:
        st.markdown("<p style='font-weight:bold; color:#a0aec0; font-size:12px;'>[ ACTIVE WORKSPACES ]</p>", unsafe_allow_html=True)
        mode_workspaces = {k: v for k, v in st.session_state.forge_workspaces.items() if v.get("type") == mode}
        
        if not mode_workspaces:
            st.info("No active projects found for this engine.")
        else:
            cols = st.columns(2)
            for idx, (ws_name, ws_data) in enumerate(mode_workspaces.items()):
                with cols[idx % 2]:
                    with st.container(border=True):
                        st.markdown(f"<span style='border:1px solid #00f3ff; color:#00f3ff; padding:2px 8px; border-radius:4px; font-size:10px;'>{ws_data['type']}</span>", unsafe_allow_html=True)
                        st.markdown(f"<h4 style='color:#ffffff; font-weight:800; margin-top:10px;'>{ws_name}</h4>", unsafe_allow_html=True)
                        st.markdown(f"<p style='font-size: 11px; color: #718096;'>ONLINE | Logs: {len(ws_data['chat'])}</p>", unsafe_allow_html=True)
                        
                        c_btn1, c_btn2 = st.columns([7, 3])
                        with c_btn1:
                            if st.button("ENTER", key=f"open_{ws_name}", use_container_width=True):
                                st.session_state.current_forge_ws = ws_name
                                st.rerun()
                        with c_btn2:
                            if st.button("DEL", key=f"del_{ws_name}", use_container_width=True):
                                del st.session_state.forge_workspaces[ws_name]
                                st.rerun()

# ==========================================
# 🖥️ ステージ3：ワークスペース内部画面
# ==========================================
else:
    ws_name = st.session_state.current_forge_ws
    ws_data = st.session_state.forge_workspaces[ws_name]
    ws_type = ws_data.get("type", "APP")
    
    if "retries" not in ws_data: ws_data["retries"] = 0
    if "media" not in ws_data: ws_data["media"] = ""
    
    col_back, col_title = st.columns([2, 8])
    with col_back:
        if st.button("← BACK TO PROJECT LIST", use_container_width=True):
            st.session_state.current_forge_ws = None
            st.rerun()
    with col_title:
        st.markdown(f"<h3 style='color:#ffffff; font-weight:800; margin-top:-5px;'>[ {ws_name} ] <span style='font-size:14px; color:#718096; font-weight:normal;'>| ENGINE: {ws_type}</span></h3>", unsafe_allow_html=True)

    st.markdown("---")

    with st.sidebar:
        st.markdown(f"<div style='text-align:center; font-weight:800; color:#00f3ff; margin-bottom:10px;'>[ {ws_name} ]</div>", unsafe_allow_html=True)
        with st.form("forge_sidebar_form", clear_on_submit=True):
            placeholder_text = "Type your prompt here..."
            if ws_type == "APP": placeholder_text = "例：シンプルな計算機アプリ"
            elif ws_type == "IMAGE": placeholder_text = "例：サイバーパンクな都市"
            elif ws_type == "VIDEO": placeholder_text = "例：コーヒーが弾ける動画"
            elif ws_type == "SLIDE": placeholder_text = "例：AIの未来について5枚"
            
            forge_prompt = st.text_area("PROMPT", placeholder=placeholder_text, height=150, label_visibility="collapsed")
            submitted = st.form_submit_button("EXECUTE", use_container_width=True, type="primary")
        
        st.markdown("<style>iframe[title*='mic'] { mix-blend-mode: multiply; opacity: 0.8; margin-top: 10px; }</style>", unsafe_allow_html=True)
        spoken_text = speech_to_text(language='ja', start_prompt="🎙️ VOICE COMMAND", stop_prompt="🛑 SEND", use_container_width=True, just_once=True, key='Forge_STT')

    col_log, col_preview = st.columns([3, 7])
    
    with col_log:
        st.markdown("<p style='font-weight:bold; color:#a0aec0; font-size:12px;'>[ SYSTEM CONSOLE ]</p>", unsafe_allow_html=True)
        core_height = 200 
        v_data = st.session_state.ai_voice_base64 if st.session_state.ai_voice_base64 else ""
        autoplay = "autoplay" if st.session_state.just_generated_audio else ""
        st.session_state.just_generated_audio = False 

        core_html = MASTER_CORE_TEMPLATE.replace("H_VAL", str(core_height)).replace("MAX_Wpx", "200").replace("V_DATA", v_data).replace("A_PLAY", autoplay)
        st.components.v1.html(core_html, height=core_height + 10)

        with st.container(height=400, border=False):
            if not ws_data["chat"]:
                st.info("System Ready. Waiting for input...")
            for m in ws_data["chat"]:
                with st.chat_message(m["role"], avatar="👤" if m["role"]=="user" else "🤖"):
                    st.markdown(m["content"])

    with col_preview:
        st.markdown("<p style='font-weight:bold; color:#a0aec0; font-size:12px;'>[ CANVAS PREVIEW ]</p>", unsafe_allow_html=True)
                
        if st.session_state.auto_fix_prompt:
            st.warning("⚙️ RECOVERY PROTOCOL INITIATED: AI is fixing the code...")
            
        elif ws_type == "APP":
            if ws_data["code"]:
                st.download_button(label="[ DOWNLOAD .py ]", data=ws_data["code"], file_name=f"{ws_name.replace(' ', '_')}.py", mime="text/plain", use_container_width=True)
                with st.container(border=True):
                    try:
                        exec(ws_data["code"], globals())
                        ws_data["retries"] = 0
                    except Exception as e:
                        st.error(f"RUNTIME ERROR:\n{e}")
                        if ws_data["retries"] < 3:
                            ws_data["retries"] += 1
                            st.session_state.auto_fix_prompt = f"実行時に以下のエラーが発生しました。修正して！\n\n【エラー内容】\n{e}"
                            st.rerun()
                        else:
                            st.error("❌ 自己修復が上限に達しました。")
                with st.expander("📝 MANUAL OVERRIDE (SOURCE CODE)"):
                    edited_code = st.text_area("Python Code", value=ws_data["code"], height=300)
                    if st.button("UPDATE CODE"):
                        ws_data["code"] = edited_code
                        st.rerun()
            else:
                st.info("Canvas is empty.")
                
        elif ws_type == "IMAGE":
            if ws_data["media"]:
                image_url = f"https://image.pollinations.ai/prompt/{urllib.parse.quote(ws_data['media'])}?width=800&height=450&nologo=true"
                st.image(image_url, use_container_width=True)
                
                with st.expander("⚙️ PROMPT TUNING (ENGLISH)"):
                    edited_media = st.text_area("Prompt", value=ws_data["media"], height=100)
                    if st.button("UPDATE IMAGE"):
                        ws_data["media"] = edited_media
                        st.rerun()
            else:
                st.info("Canvas is empty.")
                
        elif ws_type == "VIDEO":
            if ws_data["code"]:
                st.info("🔌 API Integration Pending: ここに動画生成APIを接続しMP4を表示します。")
                st.success("✅ スクリプトとビデオプロンプトが準備完了しました。")
                
                with st.expander("📝 STORYBOARD & PROMPT", expanded=True):
                    edited_code = st.text_area("Markdown Editor", value=ws_data["code"], height=300)
                    if st.button("UPDATE SCRIPT"):
                        ws_data["code"] = edited_code
                        st.rerun()
                st.markdown(ws_data["code"])
            else:
                st.info("Canvas is empty.")
                
        elif ws_type == "SLIDE":
            if ws_data["code"]:
                st.success("✅ プレゼン構成が完了しました。下のボタンから PowerPoint をダウンロードできます！")
                try:
                    prs = Presentation()
                    slides_data = json.loads(ws_data["code"])
                    
                    for slide_info in slides_data:
                        slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(slide_layout)
                        title = slide.shapes.title
                        content = slide.placeholders[1]
                        
                        title.text = slide_info.get("title", "No Title")
                        tf = content.text_frame
                        for bullet in slide_info.get("bullets", []):
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.font.size = Pt(24)
                            
                    pptx_io = io.BytesIO()
                    prs.save(pptx_io)
                    pptx_io.seek(0)
                    
                    st.download_button(
                        label="[ DOWNLOAD .pptx ]",
                        data=pptx_io,
                        file_name=f"{ws_name.replace(' ', '_')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        type="primary",
                        use_container_width=True
                    )
                except Exception as e:
                    st.error(f"パワポの生成に失敗しました: {e}")
                    st.code(ws_data["code"])
                
                with st.expander("⚙️ SLIDE DATA (JSON)"):
                    edited_code = st.text_area("JSON Editor", value=ws_data["code"], height=300)
                    if st.button("UPDATE SLIDES"):
                        ws_data["code"] = edited_code
                        st.rerun()
            else:
                st.info("Canvas is empty.")

    # ==========================================
    # AI実行ロジック（4つの脳みそ切り替え）
    # ==========================================
    trigger_prompt = forge_prompt if submitted else spoken_text if spoken_text else None
    
    is_auto_fix = False
    if st.session_state.auto_fix_prompt:
        trigger_prompt = st.session_state.auto_fix_prompt
        st.session_state.auto_fix_prompt = ""
        is_auto_fix = True

    if trigger_prompt:
        if is_auto_fix:
            ws_data["chat"].append({"role": "system", "avatar": "⚠️", "content": f"⚙️ AUTO-HEALING INITIATED:\n{trigger_prompt}"})
            sys_msg = "Auto-Healing in progress..."
        else:
            ws_data["chat"].append({"role": "user", "avatar": "👤", "content": trigger_prompt})
            ws_data["retries"] = 0
            sys_msg = f"Processing {ws_type.split(' ')[0]} Request..."

        with st.spinner(sys_msg):
            try:
                history_text = "【これまでの会話履歴】\n" + "\n".join([f"{msg['role']}: {msg['content']}" for msg in ws_data["chat"][:-1]])
                
                # 🧠 究極のシステムプロンプト（The Master Prompts）
                if "APP" in ws_type:
                    system_instruction = f"""
                    あなたは世界トップクラスのStreamlitアプリ開発者（シニアアーキテクト）です。
                    ユーザーの曖昧な指示から「本当に必要な機能」を推測し、見た目も美しく、バグのない完璧なPythonコードを提供します。
                    
                    【絶対遵守のレイアウト保護ルール】
                    1. `st.sidebar`、`st.set_page_config`、`st.chat_input` はOSの親画面を破壊するため【絶対に使用禁止】。チャットが必要な場合は `st.text_input` と `st.button` で代用せよ。
                    
                    【プロフェッショナルな開発要件】
                    1. 生成するアプリは完全独立型とし、`st.session_state` を活用して状態を適切に管理すること。
                    2. 単なる機能だけでなく、CSS（`st.markdown`）を駆使して「Neumorphism（ニューモーフィズム）」「Glassmorphism（グラスモーフィズム）」などのモダンで美しいUIを必ず実装すること。
                    3. 外部API（Gemini等）を使う場合は、必ずアプリ内に `st.text_input(..., type="password")` でキーを入力させる安全な設計にすること。
                    4. エラーハンドリング（`try-except`）を徹底し、ユーザーに優しいエラーメッセージ（`st.error` / `st.warning`）を表示すること。
                    5. コードは `# ...中略...` などの省略を絶対にせず、1行目から最後まで完全に出力すること。
                    
                    【出力フォーマット】
                    ```python
                    （ここに完全なコード）
                    ```
                    💡 次の拡張アイデア：
                    （アプリをより良くするためのプロ目線の提案を3つ）
                    
                    {history_text}
                    """
                
                elif "IMAGE" in ws_type:
                    system_instruction = f"""
                    あなたはMidjourneyやStable Diffusion、Imagen等の画像生成AIを完璧に操る、世界トップクラスの「AIプロンプトエンジニア兼アートディレクター」です。
                    ユーザーの簡単な日本語の要望から、画像生成AIが最も高品質で芸術的な画像を出力できる【究極の英語プロンプト】を構築します。
                    
                    【プロンプト構築の原則】
                    以下の要素をカンマ区切りの英語で緻密に記述すること：
                    1. Subject (主題): 構図、ポーズ、服装、表情
                    2. Medium (媒体): 写真、油絵、3Dレンダリング、水彩画、ベクターアートなど
                    3. Environment (環境): 背景、時間帯、天候、雰囲気
                    4. Lighting (照明): Cinematic lighting, volumetric lighting, rim lighting, soft softbox, neon lightsなど
                    5. Camera/Lens (カメラ設定): 35mm lens, f/1.8, macro photography, depth of field, drone shotなど
                    6. Style/Engine (スタイル): Unreal Engine 5, Octane Render, 8k resolution, highly detailed, masterpieceなど
                    
                    【出力フォーマット】
                    必ず以下の隠しタグ内に英語のプロンプトを記述すること。
                    [IMAGE_PROMPT: (ここに構築した緻密な英語プロンプト)]
                    
                    プロンプトの後に、日本語で「どのような意図でこのプロンプトを設計したか」の解説と、「さらに別のテイストにするためのアイデア」を簡潔に添えること。
                    
                    {history_text}
                    """
                    
                elif "VIDEO" in ws_type:
                    system_instruction = f"""
                    あなたはハリウッドで活躍する一流の映像ディレクター兼、SoraやVeoなどの最先端「動画生成AI」のプロンプトスペシャリストです。
                    ユーザーの要望から、プロの映像作品を作るための「絵コンテ構成」と、AIに直接入力する「英語の動画生成プロンプト」を作成します。
                    
                    【動画プロンプト（英語）の必須要素】
                    - Camera Movement (カメラワーク): Panning, Tilt, Tracking shot, Dolly zoom, FPV drone shot, Slow motionなど
                    - Scene Description (情景): 物理的な動き、光の反射、パーティクル（埃や火の粉）、被写界深度の変化
                    - Resolution/Style (画質): 4k, photorealistic, cinematic, 60fps
                    
                    【出力フォーマット（Markdown）】
                    ## 🎬 映像コンセプト
                    （どのような映像になるかの日本語解説）
                    
                    ## 🎥 シーン構成（絵コンテ）
                    - **Scene 1 (0:00-0:03):** （日本語でのシーン説明）
                    - **Scene 2 (0:03-0:06):** （日本語でのシーン説明）
                    ...
                    
                    ## 🤖 AI用動画生成プロンプト (English)
                    `（ここに動画生成AIにそのままコピペできる、すべてのシーンを統合した高品質な英語プロンプトを記述）`
                    
                    {history_text}
                    """
                    
                elif "SLIDE" in ws_type:
                    system_instruction = f"""
                    あなたはマッキンゼーなどのトップコンサルティングファームで活躍する、一流のプレゼン・ストラテジストです。
                    ユーザーのテーマに基づき、聴衆を惹きつける論理的で説得力のあるプレゼンテーション資料（スライド構成）を作成します。
                    
                    【スライド構築の絶対ルール】
                    1. 構成は「結序破急」または「PREP法」など、論理的なストーリーテリングを意識すること。
                    2. スライドとスライドの間は必ず `---` (ハイフン3つ) のみで区切ること（システムがこれでスライドを分割します）。
                    3. 1枚のスライドの情報量は多すぎず、視覚的に分かりやすいMarkdown（箇条書き、太字、引用）を使うこと。
                    
                    【出力フォーマット例】
                    # スライドタイトル1
                    ### サブタイトル
                    - ポイントA
                    - ポイントB
                    > 印象的な引用やメッセージ
                    
                    ---
                    
                    # スライドタイトル2
                    ...
                    
                    スライドを出力した後に、「このプレゼンを話す際のアドバイス（トークスクリプトのヒント）」を日本語で添えること。
                    
                    {history_text}
                    """

                # Geminiの実行
                model = genai.GenerativeModel('gemini-2.5-flash')
                response = model.generate_content(system_instruction + "\n\nボスの現在の指示: " + trigger_prompt)
                ai_text = response.text
                
                # モードごとの後処理（パース）
                reply_text = ai_text
                
                if "APP" in ws_type:
                    code_match = re.search(r'```python\n(.*?)\n```', ai_text, re.DOTALL)
                    if code_match:
                        ws_data["code"] = code_match.group(1)
                        reply_text = ai_text.replace(code_match.group(0), "").strip() or "アプリケーションのコードを構築しました。"
                
                elif "IMAGE" in ws_type:
                    prompt_match = re.search(r'\[IMAGE_PROMPT:\s*(.*?)\]', ai_text)
                    if prompt_match:
                        ws_data["media"] = prompt_match.group(1).strip()
                        reply_text = ai_text.replace(prompt_match.group(0), "").strip()
                
                elif "VIDEO" in ws_type or "SLIDE" in ws_type:
                    ws_data["code"] = ai_text 
                    reply_text = "資料の作成が完了しました。右のプレビュー画面を確認してください。"

                # 音声の生成と保存
                tts = gTTS(text=re.sub(r'[*#`_]', '', reply_text[:200]), lang='ja')
                audio_fp = io.BytesIO()
                tts.write_to_fp(audio_fp)
                st.session_state.ai_voice_base64 = base64.b64encode(audio_fp.getvalue()).decode()
                st.session_state.just_generated_audio = True
                
                ws_data["chat"].append({"role": "assistant", "avatar": "🤖", "content": reply_text})
                st.rerun()
            except Exception as e:
                st.error(f"Error: {e}")

# ==========================================
# 📥 アプリ保存モジュール (APPモード時のみ表示)
# ==========================================
if st.session_state.current_forge_ws and "APP" in ws_data.get("type", ""):
    st.markdown("---")
    st.markdown("#### 💾 SAVE TO APP ARCHIVE")
    with st.expander("📦 新しいミニアプリとしてインストール", expanded=False):
        app_filename = st.text_input("アプリのファイル名（半角英数字）", placeholder="例: my_calculator")
        app_code_input = st.text_area("保存するPythonコードを貼り付け", height=250, value=ws_data.get("code", ""))
        
        if st.button("ARCHIVE にインストール ⚡", use_container_width=True, type="primary"):
            if app_filename and app_code_input:
                safe_name = app_filename.replace(" ", "_").lower()
                if not safe_name.endswith(".py"): safe_name += ".py"
                os.makedirs("forge_apps", exist_ok=True)
                try:
                    with open(os.path.join("forge_apps", safe_name), "w", encoding="utf-8") as f:
                        f.write(app_code_input)
                    st.success(f"✅ インストール完了！ `{safe_name}` をAPP ARCHIVEに保存しました。")
                    st.balloons()
                except Exception as e:
                    st.error(f"保存エラー: {e}")